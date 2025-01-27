# -*- coding: utf-8 -*-
"""
Created on Tue Jul  9 15:19:09 2019
Letzte Änderungen
27.11.19: Belichtungszeit richtig angeben
28.11.19: Convert Problem, Index von 5 auf 4 gesetzt
3.8.20: Mit Unterordner und PP-Name nach Ordner
10.11.20 Anscheinend neue EXIF Formatierung.
24.1.21 kleiner Bugfix
17.1.23 Datum statt Kameradaten; Ortsinformationen
26.1.25 Farbiger Hintergrund / Durchschnittsfarbe

TO DO:
- schwarze Schrift Datums/Ort passend zu Hintergrund
- Ambient Effekt ?



Hilfreiche Links:


https://www.shibutan-bloomers.com/python-libraly-pptx-2_en/6310/

Chrome Webdriver:
https://googlechromelabs.github.io/chrome-for-testing/#stable

Karternmaterial 
https://leaflet-extras.github.io/leaflet-providers/preview/
"""


# fuer pptx und exifread: pip install python-pptx und pip install exifread
import os
import time

from pptx import Presentation
from pptx.util import Inches,Cm
from pptx.dml.color import RGBColor
from PIL import Image
from PIL.ExifTags import TAGS
from pptx.util import Pt
from geopy.geocoders import Nominatim
geoLoc = Nominatim(user_agent="GetLoc")
import cv2, numpy as np
from sklearn.cluster import KMeans
import folium
from selenium import webdriver # nur notwendig wenn Karten ausgegeben werden sollen


Ordner="Fotoabend2024"  # Folder containing all photos
Mit_UnterOrdner=False
Datum_Statt_Kamera_Daten=True # Datum in richtige Form bringen
Overview_Slide=False
neu_konvertieren=True # muss bei erster Ausführung true sein
karte_ausgeben=True
Karte_Neu_Berechnen=True
GeoInfos_Einbauen=True
fontsize=12


#Korrekte Umrechnung Zeiten
def belichtungszeit(zeit):
    if zeit<1:
        zeit=int(1/zeit)
        zeit="1/"+str(zeit)
    return zeit

#Minuten und Sekunden addieren
def gps_converter(North,East):
    North_Dezi=North[0]+North[1]/60.+North[2]/3600.
    East_Dezi = East[0]+ East[1]/60.+ East[2]/3600.
    return North_Dezi,East_Dezi


#Exifdaten auslesen
def get_exif(fn):
    #print("Pfad für EXIF: ",fn)
    ret = {}
    i = Image.open(fn)
    info = i._getexif()  
    keys = list(info.keys())
    keys = [k for k in keys if k in TAGS]
    for tag, value in info.items():
        decoded = TAGS.get(tag, tag)
        ret[decoded] = value
    return ret

#Bildgröße berechnen
def px_to_inches(path):    
    im = Image.open(path)
    width = im.width 
    height = im.height 
    #print("Image Info",im.info,"breite",width,"höhe",height)
    return (width, height)

#Datumsformat von Exif verbessern
def change_date_format(date):
    zwi=date.split(" ")
    datum=zwi[0].split(":")
    zeit=zwi[1]
    return datum[2]+"-"+datum[1]+"-"+datum[0]+" "+zeit


def HTML_TO_PNG(Peak):
    driver = webdriver.Chrome()
    driver.set_window_size(1000, 1000)  # choose a resolution
    driver.get("file:///Users/stephan/Documents/GitHub/Bilder_zu_PowerPoint/maps_html/" + Peak +".html") #Where to find HTML
    #driver.get("file:C:\\Users\\Stephan\\OneDrive\\PythonTruhe\\GEO_in_Bild\\maps_html\\" + Peak +".html") #Where to find HTML

    time.sleep(2) # Waiting time for page loading
    driver.save_screenshot("maps_png/"+Peak+"_largemap.png")    # Save Screenshot of HTML     
    driver.quit() #close Chrome

def visualize_colors(n_clusters, Pfad):
    # https://stackoverflow.com/questions/43111029/how-to-find-the-average-colour-of-an-image-in-python-with-opencv
    # Get the number of different clusters, create histogram, and normalize
    image = cv2.imread(Pfad)
    image = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
    reshape = image.reshape((image.shape[0] * image.shape[1], 3))
    #print("reshape",reshape)
    # Find and display most dominant colors
    cluster = KMeans(n_clusters).fit(reshape)
    
    labels = np.arange(0, len(np.unique(cluster.labels_)) + 1)
    (hist, _) = np.histogram(cluster.labels_, bins = labels)
    hist = hist.astype("float")
    hist /= hist.sum()
    # Create frequency rect and iterate through each cluster's color and percentage
    rect = np.zeros((50, 300, 3), dtype=np.uint8)
    colors = sorted([(percent, color) for (percent, color) in zip(hist,cluster.cluster_centers_)])
    start = 0
    for (percent, color) in colors:
        print(color, "{:0.2f}%".format(percent * 100))
        end = start + (percent * 300)
        cv2.rectangle(rect, (int(start), 0), (int(end), 50), \
                      color.astype("uint8").tolist(), -1)
        start = end
    return rect, colors


def zeichne_karte(North_Dezi,East_Dezi,name,start_ausschnitt=0.0025):
    karte_folium= folium.Map( zoom_start=10, scrollWheelzoom=True, zoom_control=True,  tiles="CartoDB.Positron") 
    folium.Marker((North_Dezi,East_Dezi),icon=folium.DivIcon(icon_size=("auto",20), icon_anchor=(13,13),html=f"""
                <div>
                <svg width="26" height="26">
                <circle cx="13" cy="13" r="13" fill="blue" />
                </svg>
                </div>""")).add_to(karte_folium)
    lon_min=(East_Dezi)-start_ausschnitt # left/east
    lon_max=(East_Dezi)+start_ausschnitt # right/west
    lat_min=(North_Dezi)-start_ausschnitt # up/north
    lat_max=(North_Dezi)+start_ausschnitt # down/south
    karte_folium.fit_bounds([(lat_min,lon_min),(lat_max,lon_max)])    
    karte_folium.save("maps_html/"+name+".html")


# Erstelle benötigte Ordner:
if os.path.isdir('maps_html')==False:
    os.mkdir()

if os.path.isdir('maps_png')==False:
    os.mkdir()
    
if os.path.isdir('Farbspekten_Bilder')==False:
    os.mkdir()


# Lösche alte Powerpointdatei, da nicht überschrieben wird
try:
    os.remove('Bilder_Fototron'+Ordner+'.pptx')
except:
    print("Keine Datei zum Löschen")


if Mit_UnterOrdner:
    JPG_Dateien=[]
    for sub_ordner in os.listdir(Ordner):
        print (sub_ordner)
        zwi=(os.listdir(Ordner+"/"+sub_ordner))
        for Dateien in zwi:
            JPG_Dateien.append(sub_ordner+"/"+Dateien)
    print(JPG_Dateien)

else:
    #Liste alle Dateien im Ordner aus
    JPG_Dateien=os.listdir(Ordner)
    #JPG_Dateien.sort()
    #print("Alle Dateien im Ordner:",JPG_Dateien)
    #auf Macs gibst eine ds_store Datei die vorher gelöscht werden muss
    #JPG_Dateien=[i for i in JPG_Dateien if i[-2:]!="re"] #ds_store Datei rauslöschen

print("Entferne überflüssige Dateien.")
JPG_Dateien=[i for i in JPG_Dateien if i.find("DS_Store")<0] #ds_store Datei rauslöschen
JPG_Dateien=[i for i in JPG_Dateien if i.find("Thumbs.db")<0] #Thumbs.db Datei rauslöschen
JPG_Dateien=[i for i in JPG_Dateien if i.find(".MOV")<0] #Thumbs.db Datei rauslöschen

# Konvertiere alle Dateien in JPEG, da es manche Dateien als Format MPO haben, was zu Problemen führt
if neu_konvertieren:
    for image in JPG_Dateien:
        #nur umwandeln wenn Datei noch nicht exisitert
        if (Ordner+"/"+image[:]).find("Convert")>-1:
            print(Ordner+"/"+image[:],"Convert_JPEG Dateien existiern schon")
        else:
            im = Image.open(Ordner+"/"+image)
            icc_profile = im.info.get('icc_profile')
            im.save(Ordner+"/"+image[:-4] + "_Convert.JPEG", "JPEG",icc_profile=icc_profile)


# nur die Dateien Listen die kein Convert enthalten
JPG_Dateien=[i for i in JPG_Dateien if i.find("Convert")<0] 
#JPG_Dateien=sorted(JPG_Dateien, key=lambda item: (int(item.partition('.')[0]) if item[0].isdigit() else float('inf'), item))
print("Dateien zu Bearbeiten:",JPG_Dateien)

Bild_Nummer=1
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_size = (33.867 , 19.05) # in cm. Größe Folien
prs.slide_width, prs.slide_height = Cm(slide_size[0]), Cm(slide_size[1])

for index,Datei in enumerate(JPG_Dateien[0:5]):
    print("Index",index)
    blank_slide_layout = prs.slide_layouts[6]
    print("Datei",Datei)
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = right = bottom = Cm(0.5) # Platzierung der Bilder
    pfad_convert=Ordner+"/"+Datei[:-4]+"_Convert.JPEG"
    print("JPEG in Powerpoint",pfad_convert)

    # Bild hinzufügen
    img = px_to_inches(pfad_convert)
    Bildhoehe=slide_size[1]-2
    Bild_ratio=img[0]/img[1]*Bildhoehe
    Folienmitte=slide_size[0]/2.  
    left=Folienmitte-Bild_ratio/2.  
    print("left",left,"Folienmitte",Folienmitte,"Bild_Ratio",Bild_ratio)
    pic = slide.shapes.add_picture(pfad_convert, Cm(left), top, height= Cm(Bildhoehe)) #Maße der Bilder
 
 
 
    #Durchschnittsfarbe bestimmen und als Hintergrund setzen:
    visualize,colors = visualize_colors(5,pfad_convert)
    visualize = cv2.cvtColor(visualize, cv2.COLOR_RGB2BGR)
    #print(colors)
    #cv2.imshow('visualize', visualize)
    cv2.imwrite("Farbspekten_Bilder/"+str(index)+"_farben.jpeg", visualize)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(int(colors[-1][1][0]), int(colors[-1][1][1]), int(colors[-1][1][2]))

    
    #Platzierung der EXIF-Daten
    left = top = width = height = Inches(1)
    left = Cm(1.2)
    top  = Cm(slide_size[1]-1.5) #0.1 oben
    width= Cm(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    exif_text = ""
    p = tf.paragraphs[0]
    run_kamera = p.add_run()

    top  = Cm(slide_size[1]-1.0) #0.1 oben
    locationBox = slide.shapes.add_textbox(left, top, width, height)
    tf_location = locationBox.text_frame
    q = tf_location.paragraphs[0]
    run_location = q.add_run()

    if Datei[-3:]!="png":
        #if True:
        try:
            ausgabe=get_exif(Ordner+"/"+Datei)
            belichtungszeit_wert=belichtungszeit( ausgabe["ExposureTime"]  )

            #print("Ausgabe",ausgabe)
            print("Brennweite",ausgabe["FocalLength"])
            print("ISO",ausgabe["ISOSpeedRatings"])
            #print("FNumber/Blende",ausgabe["FNumber"],round( ausgabe["FNumber"][0]/ausgabe["FNumber"][1],1))
            print("Exposure Time",ausgabe["ExposureTime"],str(round(ausgabe["ExposureTime"],6)))
            print("Model",ausgabe["Model"])
            #print("Lens Model",ausgabe["LensModel"])
            exif_text=exif_text+str(Bild_Nummer)+")   " 
            exif_text=exif_text+"Brennweite: "+str(ausgabe["FocalLength"]) + "  ISO: "+str(ausgabe["ISOSpeedRatings"]) 
            try:
                # muss hier noch verbessert werden
                exif_text=exif_text+"  Blende: "+str(round( ausgabe["ApertureValue"]))
                exif_text=exif_text+"  Blende: "+str(round( ausgabe["MaxApertureValue"]))
            except:
                print("Problem mit Blende")
            exif_text=exif_text+"  Belichtungszeit: "+str(belichtungszeit_wert)+" s"
            
            if Datum_Statt_Kamera_Daten:
                run_kamera.text=change_date_format(ausgabe["DateTime"])
            else:
                run_kamera.text=exif_text

        except Exception as fehler:
            print("Fehler bei EXIF Daten",fehler)
            

        # Geo Informationen einbauen    
        try: 
            if GeoInfos_Einbauen:
                North_Dezi,East_Dezi=gps_converter(ausgabe["GPSInfo"][2],ausgabe["GPSInfo"][4])
                
                # passing the coordinates
                locname = geoLoc.reverse(f"{North_Dezi}, {East_Dezi}")

                # printing the address/location name
                print("locname.address",locname.address)
                run_location.text=locname.address
                if Karte_Neu_Berechnen:
                    zeichne_karte(North_Dezi,East_Dezi,str(index))
                    HTML_TO_PNG(str(index))
                pic_map = slide.shapes.add_picture("maps_png/"+str(index)+"_largemap.png", Cm(27.8), Cm(13.82), height= Cm(5)) #Maße der Bilder
            Bild_Nummer=Bild_Nummer+1
            print("-----------------------------------")
        except Exception as fehler:
            print("Fehler bei Exifdaten:",fehler)
            #exif_text=Datei
            print("-----------------------------------")

    font = run_kamera.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)
    font = run_location.font
    font.name = 'Calibri'
    font.size = Pt(fontsize)




# Overview
if Overview_Slide:
    # Empty Page for Overview
    blank_slide_layout = prs.slide_layouts[6]
    print("Datei",Datei)
    slide = prs.slides.add_slide(blank_slide_layout)
    index=0
    for hoehe in range(4):
        for laenge in range(5):
            try:
                left = top = right = bottom = Inches(0.5) # Platzierung der Bilder
                left = Inches(0.5)+Inches(laenge*2)
                top =  Inches(0.5)+Inches(hoehe*2)
                pfad_convert=Ordner+"/"+JPG_Dateien[index][:-4]+"_Convert.JPEG"
                print("JPEG in Powerpoint Mehrfachplot",pfad_convert)
                pic = slide.shapes.add_picture(pfad_convert, left, top, height= Inches(1.5)) #Maße der Bilder
        
                index=index+1
            except:
                print("Keine Bilder mehr für übersicht")



prs.save('Prensentation_'+Ordner+'.pptx')
